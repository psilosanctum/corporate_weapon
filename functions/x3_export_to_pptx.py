from pptx import Presentation
from pptx.util import Inches  #Can simply transfer number to inches

def createPPTX():
    try:
        print('Creating PPTX from template...')
        prs = Presentation('template.pptx')
        bar_transactions = 'figures/bar_transactions.png'
        line_transactions = 'figures/line_transactions.png'

        page_id = {i+1: slide.slide_id for i, slide in enumerate(prs.slides)}
        slide1 = prs.slides.get(page_id[1])  #Get the first slide
        pic = slide1.shapes.add_picture(bar_transactions, 
                                        left=Inches(0.5),
                                        top=Inches(1.25),
                                        width=Inches(12.5),
                                        height=Inches(6))

        slide2 = prs.slides.get(page_id[2])  #Get the seceond slide
        pic = slide2.shapes.add_picture(line_transactions,
                                    left=Inches(0.5),
                                    top=Inches(1.25),
                                    width=Inches(12.5),
                                    height=Inches(6))

        prs.save('monthly_transaction_report.pptx')
        print('Monthly PPTX saved!')

    except:
        print("Monthly PPTX not saved!")

# createPPTX()